import io
import re
from datetime import datetime

import matplotlib
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

try:
    import arabic_reshaper
except Exception:
    arabic_reshaper = None

try:
    from bidi.algorithm import get_display
except Exception:
    get_display = None


matplotlib.use("Agg")
matplotlib.rcParams["font.family"] = ["Arial", "Tahoma", "DejaVu Sans", "sans-serif"]


COR = RGBColor(0xE8, 0x62, 0x4A)
NAVY = RGBColor(0x2D, 0x3A, 0x52)
GRAY = RGBColor(0xF5, 0xF6, 0xF8)
LIGHT_GRAY = RGBColor(0xEC, 0xEE, 0xF3)
MUTED = RGBColor(0x69, 0x73, 0x85)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2D, 0x7A, 0x4F)
PALETTE = ["#E8624A", "#2D3A52", "#A0B0C8", "#F4A990", "#8B9DB8", "#F07A65"]
SLIDE_MARGIN_X = 0.18
CONTENT_GAP_X = 0.22
TWO_COL_W = 6.35


def _contains_arabic(text):
    return bool(re.search(r"[\u0600-\u06FF]", str(text)))


def _shape_text(text):
    text = str(text)
    if not text:
        return ""
    if _contains_arabic(text) and arabic_reshaper and get_display:
        try:
            return get_display(arabic_reshaper.reshape(text))
        except Exception:
            return text
    return text


def _display_label(text, max_len=None):
    text = str(text)
    if max_len and len(text) > max_len:
        text = text[: max_len - 1].rstrip() + "…"
    return _shape_text(text)


def _fig_bytes(fig):
    buffer = io.BytesIO()
    fig.savefig(buffer, format="png", dpi=130, bbox_inches="tight", facecolor="white", edgecolor="none")
    plt.close(fig)
    buffer.seek(0)
    return buffer


def _pie(labels, values, title):
    fig, ax = plt.subplots(figsize=(5.4, 4.2))
    fig.patch.set_facecolor("white")
    wedges, _, autotexts = ax.pie(
        values,
        autopct="%1.1f%%",
        colors=PALETTE[: len(values)],
        startangle=90,
        pctdistance=0.78,
        wedgeprops=dict(linewidth=2, edgecolor="white"),
    )
    for item in autotexts:
        item.set_fontsize(9)
        item.set_color("white")
        item.set_fontweight("bold")
    ax.legend(
        wedges,
        [_display_label(label, 28) for label in labels],
        loc="lower center",
        bbox_to_anchor=(0.5, -0.15),
        ncol=2,
        fontsize=8,
        frameon=False,
    )
    ax.set_title(_shape_text(title), fontsize=11, fontweight="bold", color="#2D3A52", pad=10)
    plt.tight_layout()
    return _fig_bytes(fig)


def _bar_h(labels, values, title, color="#E8624A"):
    fig, ax = plt.subplots(figsize=(7, max(3.5, len(labels) * 0.45)))
    fig.patch.set_facecolor("white")
    ax.barh(range(len(labels)), values, color=color, edgecolor="white", linewidth=0.5, height=0.65)
    ax.set_yticks(range(len(labels)))
    ax.set_yticklabels([_display_label(label, 35) for label in labels], fontsize=8.5)
    ax.invert_yaxis()
    ax.set_title(_shape_text(title), fontsize=11, fontweight="bold", color="#2D3A52", pad=10)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_visible(False)
    ax.xaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f"{x:,.0f}"))
    ax.tick_params(axis="y", pad=6)
    max_value = max(values) if values else 1
    for idx, value in enumerate(values):
        ax.text(value + max_value * 0.01, idx, f"{value:,.0f}", va="center", fontsize=8, color="#697385")
    plt.tight_layout()
    return _fig_bytes(fig)


def _bar_v(labels, values, title):
    fig, ax = plt.subplots(figsize=(5.5, 3.5))
    fig.patch.set_facecolor("white")
    bars = ax.bar(range(len(labels)), values, color=PALETTE[: len(labels)], edgecolor="white", linewidth=0.5, width=0.55)
    ax.set_xticks(range(len(labels)))
    ax.set_xticklabels([_display_label(label, 18) for label in labels], fontsize=9)
    ax.set_title(_shape_text(title), fontsize=11, fontweight="bold", color="#2D3A52", pad=10)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f"{x:,.0f}"))
    for bar, value in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width() / 2, value * 1.01, f"{value:,.0f}", ha="center", fontsize=8, color="#697385")
    plt.tight_layout()
    return _fig_bytes(fig)


def _line_chart(labels, values, title):
    fig, ax = plt.subplots(figsize=(8, 3.2))
    fig.patch.set_facecolor("white")
    ax.plot(labels, values, color="#E8624A", marker="o", linewidth=2, markersize=5)
    ax.fill_between(range(len(values)), values, alpha=0.15, color="#E8624A")
    ax.set_xticks(range(len(labels)))
    ax.set_xticklabels(labels, rotation=45, ha="right", fontsize=8)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f"{x:,.0f}"))
    ax.set_title(_shape_text(title), fontsize=11, fontweight="bold", color="#2D3A52", pad=10)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.tick_params(axis="y", labelsize=8)
    plt.tight_layout()
    return _fig_bytes(fig)


def _add_text(slide, text, x, y, w, h, size=12, bold=False, color=None, align=PP_ALIGN.LEFT, italic=False):
    text_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = _shape_text(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return text_box


def _add_img(slide, image_buffer, x, y, w, h=None):
    if h:
        slide.shapes.add_picture(image_buffer, Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        slide.shapes.add_picture(image_buffer, Inches(x), Inches(y), width=Inches(w))


def _rect(slide, x, y, w, h, fill_color):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _table(slide, headers, rows, x, y, w, col_widths, header_color=None, row_h=0.28):
    header_color = header_color or COR
    table = slide.shapes.add_table(
        len(rows) + 1,
        len(headers),
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(row_h * (len(rows) + 1)),
    ).table

    for idx, col_width in enumerate(col_widths):
        table.columns[idx].width = Inches(col_width)

    for idx, header in enumerate(headers):
        cell = table.cell(0, idx)
        cell.text = _shape_text(header)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
        run.font.color.rgb = WHITE
        run.font.bold = True
        run.font.size = Pt(8.5)

    for row_idx, row in enumerate(rows):
        bg = GRAY if row_idx % 2 else WHITE
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = _shape_text(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
            run.font.size = Pt(8)
            run.font.color.rgb = NAVY


def build_consumption_pptx(
    df,
    col_member=None,
    col_amount=None,
    col_provider=None,
    col_ptype=None,
    col_diag=None,
    col_date=None,
    col_incident=None,
    col_chronic=None,
    col_card=None,
    col_relation=None,
    col_category=None,
    annual_premium=0,
    ibnr_pct=10,
    insurer_name="",
    client_name="",
    preparer_name="Ahmed Soufy",
    preparer_title="Corporate Medical Account Manager",
):
    df = df.copy()
    df["_amt"] = pd.to_numeric(df[col_amount], errors="coerce").fillna(0) if col_amount else 0.0

    total_paid = df["_amt"].sum()
    total_claims = len(df)
    active_members = int(df[col_member].nunique()) if col_member else 0
    avg_member = total_paid / active_members if active_members else 0
    lr_net = (total_paid / annual_premium * 100) if annual_premium else 0
    loss_ratio = (total_paid * (1 + ibnr_pct / 100) / annual_premium * 100) if annual_premium else 0
    lr_ok = loss_ratio <= 100
    now_str = datetime.now().strftime("%d/%m/%Y")
    insurer_name = insurer_name.strip() or "Insurance Co."
    client_name = client_name.strip() or "Client"

    def group_amount(column, sort="total"):
        return (
            df.groupby(column)["_amt"]
            .agg(total="sum", count="count")
            .reset_index()
            .sort_values(sort, ascending=False)
        )

    top_members = []
    top_providers = []
    provider_types = []
    diagnoses = []
    incident_rows = []
    chronic_rows = []

    if col_member:
        grouped = group_amount(col_member).head(10)
        if col_card:
            card_map = df.groupby(col_member)[col_card].first().to_dict()
            grouped["card"] = grouped[col_member].map(card_map)
        for _, row in grouped.iterrows():
            top_members.append([
                str(row.get("card", "")) if pd.notna(row.get("card")) else "",
                str(row[col_member]),
                int(row["count"]),
                f"{row['total']:,.2f}",
            ])

    if col_provider:
        for _, row in group_amount(col_provider, "count").head(10).iterrows():
            top_providers.append([str(row[col_provider])[:40], int(row["count"]), f"{row['total']:,.2f}"])

    if col_ptype:
        for _, row in group_amount(col_ptype).iterrows():
            provider_types.append([str(row[col_ptype]), int(row["count"]), f"{row['total']:,.2f}"])

    if col_diag:
        for _, row in group_amount(col_diag, "count").head(10).iterrows():
            diagnoses.append([str(row[col_diag])[:40], int(row["count"]), f"{row['total']:,.2f}"])

    summary_col = col_incident or col_category
    summary_title = "Incident Type" if col_incident else ("Service Category" if col_category else "Breakdown")
    if summary_col:
        grouped = df.groupby(summary_col)["_amt"].agg(total="sum", cnt="count").reset_index()
        for _, row in grouped.iterrows():
            incident_rows.append([str(row[summary_col]), int(row["cnt"]), f"{row['total']:,.2f}"])

    if col_chronic:
        grouped = df.groupby(col_chronic)["_amt"].agg(total="sum", cnt="count").reset_index()
        if col_member:
            member_counts = df.groupby(col_chronic)[col_member].nunique().reset_index()
            member_counts.columns = [col_chronic, "members"]
            grouped = grouped.merge(member_counts, on=col_chronic)
        else:
            grouped["members"] = grouped["cnt"]
        grouped["label"] = grouped[col_chronic].map({"Y": "Chronic", "N": "Non-Chronic"}).fillna(grouped[col_chronic])
        grouped["avg"] = grouped["total"] / grouped["members"]
        for _, row in grouped.iterrows():
            chronic_rows.append([str(row["label"]), int(row["members"]), f"{row['total']:,.2f}", f"{row['avg']:,.2f}"])

    monthly_chart = None
    if col_date:
        try:
            df["_date_parsed"] = pd.to_datetime(df[col_date], errors="coerce")
            monthly = (
                df.dropna(subset=["_date_parsed"])
                .groupby(df["_date_parsed"].dt.to_period("M"))["_amt"]
                .sum()
                .reset_index()
            )
            monthly.columns = ["Month", "Amount"]
            monthly["Month"] = monthly["Month"].astype(str)
            if not monthly.empty:
                monthly_chart = _line_chart(monthly["Month"].tolist(), monthly["Amount"].tolist(), "Monthly Consumption Trend")
        except Exception:
            monthly_chart = None

    member_chart_labels = []
    for row in top_members:
        card_no = str(row[0]).strip()
        member_name = str(row[1]).strip()
        if member_name and card_no:
            member_chart_labels.append(f"{member_name} - {card_no}")
        else:
            member_chart_labels.append(member_name or card_no)

    members_chart = _bar_h(
        member_chart_labels,
        [float(row[3].replace(",", "")) for row in top_members],
        "Top Members by Amount",
    ) if top_members else None
    members_count_chart = _bar_h(
        member_chart_labels[:5],
        [row[2] for row in top_members[:5]],
        "Top Members by Claim Count",
        color="#A0B0C8",
    ) if top_members else None
    providers_chart = _bar_h(
        [row[0] for row in top_providers[:8]],
        [float(row[2].replace(",", "")) for row in top_providers[:8]],
        "Top Providers by Amount",
        color="#2D3A52",
    ) if top_providers else None
    provider_type_chart = _bar_h(
        [row[0] for row in provider_types[:8]],
        [float(row[2].replace(",", "")) for row in provider_types[:8]],
        "Provider Type Distribution",
        color="#2D3A52",
    ) if provider_types else None
    incident_chart = _pie([row[0] for row in incident_rows], [row[1] for row in incident_rows], summary_title) if incident_rows else None
    chronic_chart = _bar_v([row[0] for row in chronic_rows], [float(row[2].replace(",", "")) for row in chronic_rows], "Chronic vs Non-Chronic") if chronic_rows else None

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def new_slide():
        return prs.slides.add_slide(blank)

    def add_header(slide, page_title):
        _rect(slide, 0, 0, 0.06, 7.5, COR)
        _rect(slide, 0.06, 0, 13.27, 0.55, NAVY)
        _add_text(slide, "Service Overview", 0.2, 0, 4, 0.55, size=13, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
        _add_text(slide, f"{client_name} | {insurer_name}", 4, 0, 5, 0.55, size=9, color=RGBColor(0xB0, 0xC0, 0xD0), align=PP_ALIGN.CENTER)
        _add_text(slide, page_title, 9.5, 0, 3.7, 0.38, size=12, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
        _add_text(slide, now_str, 9.5, 0.33, 3.7, 0.22, size=8, color=RGBColor(0x90, 0xA0, 0xB0), align=PP_ALIGN.RIGHT)

    def kpi_card(slide, x, y, w, h, label, value, highlight=False):
        _rect(slide, x, y, w, h, LIGHT_GRAY)
        _add_text(slide, label, x + 0.08, y + 0.06, w - 0.16, 0.22, size=7.5, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
        _add_text(slide, value, x + 0.08, y + 0.28, w - 0.16, 0.38, size=19, bold=True, color=COR if highlight else NAVY, align=PP_ALIGN.CENTER)

    s1 = new_slide()
    _rect(s1, 0.0, 0, 5.8, 7.5, NAVY)
    _rect(s1, 5.5, 0, 0.3, 7.5, COR)
    _rect(s1, 5.8, 0, 7.53, 7.5, GRAY)
    _add_text(s1, "CR Analysis", 0.4, 1.5, 5.0, 0.8, size=38, bold=True, color=WHITE)
    _add_text(s1, client_name, 0.4, 2.35, 5.0, 0.5, size=22, bold=True, color=COR)
    _add_text(s1, "Consumption Report", 0.4, 2.85, 5.0, 0.4, size=16, color=RGBColor(0x90, 0xA8, 0xC0))

    meta_rows = [
        ("Insurance Company", insurer_name),
        ("Issued Date", now_str),
        ("Total Claims", f"{total_claims:,}"),
        ("Active Members", f"{active_members:,}"),
    ]
    for idx, (label, value) in enumerate(meta_rows):
        box_y = 1.1 + idx * 1.32
        _rect(s1, 6.3, box_y, 6.5, 1.1, WHITE)
        _add_text(s1, label, 6.45, box_y + 0.1, 6.2, 0.25, size=8, bold=True, color=MUTED)
        _add_text(s1, value, 6.45, box_y + 0.38, 6.2, 0.45, size=16, bold=True, color=NAVY)

    _add_text(s1, f"Prepared by: {preparer_name} | {preparer_title}", 1.296, 7.16, 12.0, 0.3, size=8, color=RGBColor(0x60, 0x70, 0x80), align=PP_ALIGN.CENTER)

    s2 = new_slide()
    add_header(s2, "Members Analysis")
    kpis = [
        ("Total Paid (EGP)", f"{total_paid:,.0f}", False),
        ("Total Claims", f"{total_claims:,}", False),
        ("Active Members", f"{active_members:,}", False),
        ("Avg / Member", f"{avg_member:,.0f}", False),
        (f"Loss Ratio + IBNR {ibnr_pct:.0f}%", f"{loss_ratio:.1f}%", not lr_ok),
    ]
    for idx, (label, value, highlight) in enumerate(kpis):
        kpi_card(s2, [0.18, 2.77, 5.36, 7.95, 10.54][idx], 0.65, 2.44, 0.75, label, value, highlight)

    if top_members:
        _add_text(s2, "Top 10 Members by Total Net Payable", SLIDE_MARGIN_X, 1.52, TWO_COL_W, 0.28, size=11, bold=True, color=NAVY)
        _table(
            s2,
            ["Card No", "Member Name", "Claim Count", "Total Amount Paid"],
            top_members,
            SLIDE_MARGIN_X,
            1.84,
            TWO_COL_W,
            [1.05, 3.05, 0.95, 1.3],
        )
    if members_count_chart:
        _add_img(s2, members_count_chart, SLIDE_MARGIN_X, 5.12, TWO_COL_W, 1.86)
    if members_chart:
        chart_x = SLIDE_MARGIN_X + TWO_COL_W + CONTENT_GAP_X
        _add_text(s2, "Top Members Distribution", chart_x, 1.52, TWO_COL_W, 0.28, size=11, bold=True, color=NAVY)
        _add_img(s2, members_chart, chart_x, 1.84, TWO_COL_W, 5.14)

    s3 = new_slide()
    add_header(s3, "Billing Providers")
    if top_providers:
        _add_text(s3, "Top 10 Billing Providers by Claims", SLIDE_MARGIN_X, 0.65, TWO_COL_W, 0.28, size=11, bold=True, color=NAVY)
        _table(
            s3,
            ["Billing Provider", "Claims Count", "Total Amount Paid"],
            top_providers,
            SLIDE_MARGIN_X,
            0.97,
            TWO_COL_W,
            [3.85, 1.0, 1.5],
        )
    if providers_chart:
        _add_text(s3, "Provider Spend Distribution", SLIDE_MARGIN_X, 4.18, TWO_COL_W, 0.28, size=11, bold=True, color=NAVY)
        _add_img(s3, providers_chart, SLIDE_MARGIN_X, 4.5, TWO_COL_W, 2.48)
    if provider_types:
        right_x = SLIDE_MARGIN_X + TWO_COL_W + CONTENT_GAP_X
        _add_text(s3, "Distribution by Provider Type", right_x, 0.65, TWO_COL_W, 0.28, size=11, bold=True, color=NAVY)
        _table(s3, ["Type", "Count", "Amount"], provider_types[:7], right_x, 0.97, TWO_COL_W, [3.0, 0.95, 2.4])
    if provider_type_chart:
        right_x = SLIDE_MARGIN_X + TWO_COL_W + CONTENT_GAP_X
        _add_text(s3, "Provider Type Spend", right_x, 3.95, TWO_COL_W, 0.28, size=11, bold=True, color=NAVY)
        _add_img(s3, provider_type_chart, right_x, 4.27, TWO_COL_W, 2.71)

    s4 = new_slide()
    add_header(s4, "Diagnoses and Breakdown")
    if diagnoses:
        _add_text(s4, "Top 10 Diagnoses by Frequency", 0.18, 0.65, 8.0, 0.28, size=11, bold=True, color=NAVY)
        _table(s4, ["Diagnosis", "Count", "Total Amount Paid"], diagnoses, 0.18, 0.97, 8.0, [5.5, 1.0, 1.5])
    if incident_rows:
        _add_text(s4, f"{summary_title} Summary", 8.5, 0.65, 4.7, 0.28, size=11, bold=True, color=NAVY)
        total_count = sum(row[1] for row in incident_rows)
        total_amount = sum(float(row[2].replace(",", "")) for row in incident_rows)
        table_rows = []
        for row in incident_rows:
            amount_value = float(row[2].replace(",", ""))
            table_rows.append([
                row[0],
                row[1],
                row[2],
                f"{(row[1] / total_count * 100) if total_count else 0:.1f}%",
                f"{(amount_value / total_amount * 100) if total_amount else 0:.1f}%",
            ])
        _table(s4, ["Type", "Count", "Total", "% Count", "% Amt"], table_rows, 8.5, 0.97, 4.62, [1.5, 0.7, 1.0, 0.71, 0.71])
    if incident_chart:
        _add_img(s4, incident_chart, 8.32, 0.97, 4.88, 5.94)
    if monthly_chart:
        _add_img(s4, monthly_chart, 0.18, 4.15, 8.0, 3.2)

    s5 = new_slide()
    add_header(s5, "Cohort Analysis and Loss Ratio")
    if chronic_rows:
        _add_text(s5, "Chronic vs Non-Chronic Cohort", 0.18, 0.65, 7.0, 0.28, size=11, bold=True, color=NAVY)
        _table(s5, ["Cohort", "Active Users", "Total Amount", "Avg / Member"], chronic_rows, 0.18, 1.24, 7.0, [2.5, 1.5, 1.5, 1.5])
    if chronic_chart:
        _add_img(s5, chronic_chart, 0.2, 2.73, 6.5, 4.07)

    lr_color = GREEN if lr_ok else COR
    _rect(s5, 8.5, 0.65, 4.62, 1.1, NAVY)
    _add_text(s5, "Loss Ratio Calculation", 8.5, 0.65, 4.62, 0.4, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_text(s5, f"{loss_ratio:.1f}%", 8.5, 1.05, 4.62, 0.7, size=34, bold=True, color=lr_color, align=PP_ALIGN.CENTER)

    badge_bg = RGBColor(0xE8, 0xF5, 0xE9) if lr_ok else RGBColor(0xFF, 0xEB, 0xEE)
    _rect(s5, 8.5, 1.82, 4.62, 0.38, badge_bg)
    _add_text(s5, "Within Range" if lr_ok else "Above 100%", 8.63, 1.86, 4.32, 0.25, size=9, bold=True, color=lr_color, align=PP_ALIGN.CENTER)

    details = [
        f"Paid Claims YTD: EGP {total_paid:,.0f}",
        f"Annual Premium: EGP {annual_premium:,.0f}",
        f"Loss Ratio (net): {lr_net:.1f}%",
        f"Loss Ratio + IBNR ({ibnr_pct:.0f}%): {loss_ratio:.1f}%",
    ]
    for idx, detail in enumerate(details):
        y = 2.28 + idx * 0.48
        _rect(s5, 8.5, y, 4.62, 0.45, WHITE if idx % 2 == 0 else GRAY)
        _add_text(s5, detail, 8.65, y, 4.3, 0.45, size=9.5, bold=(idx == 3), color=lr_color if idx == 3 else NAVY)

    s6 = new_slide()
    _rect(s6, 0.0, 0.0, 13.33, 7.5, NAVY)
    _rect(s6, 0.0, 0.0, 13.33, 0.08, COR)
    _rect(s6, 0.0, 7.42, 13.33, 0.08, COR)
    _add_text(s6, "Prepared for you by:", 0.0, 2.2, 13.33, 0.4, size=11, italic=True, color=RGBColor(0x80, 0x90, 0xA0), align=PP_ALIGN.CENTER)
    _rect(s6, 5.0, 2.75, 3.3, 0.04, COR)
    _add_text(s6, preparer_name, 0.0, 2.9, 13.33, 1.1, size=52, bold=True, italic=True, color=COR, align=PP_ALIGN.CENTER)
    _add_text(s6, preparer_title, 0.0, 4.1, 13.33, 0.5, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _rect(s6, 5.5, 4.75, 2.3, 0.04, RGBColor(0x40, 0x50, 0x60))
    _add_text(s6, f"Report generated on {now_str}", 0.0, 4.9, 13.33, 0.3, size=9, color=RGBColor(0x60, 0x70, 0x80), align=PP_ALIGN.CENTER)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()
